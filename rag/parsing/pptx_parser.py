"""PPT/PPTX 课件解析模块。"""

from __future__ import annotations

import io
import os
from typing import Dict, List, Optional, Union

from PIL import Image

from ..processing.cleaners import clean_text
from ..config import ExtractionConfig
from .ocr_backends import create_ocr_backend


def _shape_text(shape) -> List[str]:
    parts: List[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text:
        parts.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    if getattr(shape, "has_chart", False):
        chart = shape.chart
        if getattr(chart, "chart_title", None) and chart.chart_title.has_text_frame:
            parts.append(chart.chart_title.text_frame.text)
    return parts


def _notes_text(slide) -> str:
    try:
        notes = slide.notes_slide
        frame = notes.notes_text_frame
        return frame.text if frame else ""
    except Exception:
        return ""


def _ocr_picture_shape(shape, backend) -> str:
    if backend is None:
        return ""
    try:
        image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
        return backend.image_to_text(image)
    except Exception:
        return ""


def extract_pages_from_pptx(
    file_path: str,
    config: Optional[Union[ExtractionConfig, Dict]] = None,
) -> List[Dict]:
    """
    从 PPTX 提取每页幻灯片文本、备注、表格文字，并尽量 OCR 图片文字。

    旧版 .ppt 需要先用 Office/WPS/LibreOffice 转成 .pptx；解析失败会抛出明确异常。
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPT 文件不存在: {file_path}")

    cfg = ExtractionConfig.from_dict(config) if isinstance(config, dict) else (config or ExtractionConfig())
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx：pip install python-pptx") from exc

    ocr_backend = None
    if cfg.enable_ocr and cfg.ocr_backend.lower() != "none":
        try:
            ocr_backend = create_ocr_backend(cfg)
        except Exception:
            ocr_backend = None

    prs = Presentation(file_path)
    pages: List[Dict] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        picture_ocr_count = 0
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))
            if getattr(shape, "shape_type", None) is not None and hasattr(shape, "image"):
                ocr_text = _ocr_picture_shape(shape, ocr_backend)
                if ocr_text.strip():
                    picture_ocr_count += 1
                    parts.append(ocr_text)

        note = _notes_text(slide)
        if note.strip():
            parts.append(note)

        text = clean_text("\n".join(parts))
        pages.append({
            "page_number": idx,
            "page": idx,
            "text": text,
            "extraction_method": "pptx" if picture_ocr_count == 0 else "pptx+image_ocr",
            "method": "pptx" if picture_ocr_count == 0 else "pptx+image_ocr",
            "char_count": len(text),
            "text_length": len(text),
            "is_ocr": picture_ocr_count > 0,
            "is_ocr_text": picture_ocr_count > 0,
            "ocr_backend": "",
            "ocr_confidence": None,
            "is_scanned_page": False,
            "failure_reason": "",
            "quality_stats": {},
            "image_ocr_count": picture_ocr_count,
        })
    return [p for p in pages if p.get("text") or p.get("image_ocr_count")]


def extract_pptx_text(
    file_path: str,
    config: Optional[Union[ExtractionConfig, Dict]] = None,
) -> str:
    """从 PPTX 提取全文。"""
    return "\n\n".join(p["text"] for p in extract_pages_from_pptx(file_path, config) if p.get("text"))
