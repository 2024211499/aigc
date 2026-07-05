# -*- coding: utf-8 -*-
"""
RAG 检索桥接层
将 1号队员的解析/检索模块封装为后端服务可调用的接口
兼容 1号代码未整合时的降级模式
"""

import os
import sys
import logging
import time
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple

from ..core.config import settings

logger = logging.getLogger(__name__)

# ─── 尝试导入 1号模块 ────────────────────────────────────────
_rag_available = False
_pdf_extractor = None
_embedding_provider = None

def _try_import_rag():
    global _rag_available
    rag_dir = settings.RAG_MODULE_DIR
    # 将 rag/ 的上级目录（项目根）加入 sys.path，以包的形式导入
    rag_parent = str(Path(rag_dir).parent)
    if os.path.isdir(rag_dir) and rag_parent not in sys.path:
        sys.path.insert(0, rag_parent)
    try:
        # 以包形式导入 1号的核心模块（rag 目录需有 __init__.py）
        from rag.retrieval.embeddings import EmbeddingProvider  # type: ignore
        from rag.pipeline.index_manager import IndexManager     # type: ignore
        _rag_available = True
        logger.info("✅ 1号 RAG 模块加载成功")
    except ImportError as e:
        logger.warning(f"⚠️  1号 RAG 模块未找到，使用内置简化检索: {e}")
        _rag_available = False

_try_import_rag()


# ════════════════════════════════════════════════════════════
#  文档解析服务
# ════════════════════════════════════════════════════════════

class DocumentParseService:
    """
    文档解析服务
    优先调用 1号模块；若不可用则使用内置简化版
    """

    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析文档，返回结构化结果

        Returns:
            {
                text: str,
                pages: int,
                chapters: list,
                knowledge_points: list,
                formulas: list,
                examples: list,
                exercises: list,
                from_ocr: bool,
                error: str
            }
        """
        ext = Path(file_path).suffix.lower()
        start = time.time()

        try:
            if _rag_available:
                result = self._parse_via_rag(file_path, ext)
            else:
                result = self._parse_builtin(file_path, ext)
            return self._validate_parse_result(result, ext)
        except Exception as e:
            logger.error(f"文档解析失败: {e}", exc_info=True)
            return self._empty_result(error=str(e))
        finally:
            elapsed = round(time.time() - start, 2)
            logger.info(f"解析耗时: {elapsed}s  文件: {Path(file_path).name}")

    def _parse_via_rag(self, file_path: str, ext: str) -> Dict[str, Any]:
        """调用 1号模块解析，优先保留页码和 OCR 元数据。"""
        from rag.config import ExtractionConfig  # type: ignore
        from rag.parsing.document_parser import parse_document_result  # type: ignore

        ocr_backend = os.getenv("PDF_OCR_BACKEND", "auto")
        online_ocr = bool(
            os.getenv("MULTIMODAL_API_KEY")
            or os.getenv("OPENAI_API_KEY")
            or os.getenv("BAIDU_OCR_API_KEY")
            or os.getenv("TENCENT_SECRET_ID")
            or os.getenv("VIVO_OCR_APP_ID")
        )
        max_pages_raw = os.getenv("PDF_OCR_MAX_PAGES", "").strip()
        max_pages = int(max_pages_raw) if max_pages_raw.isdigit() else None
        cfg = ExtractionConfig(
            enable_ocr=True,
            enable_online_ocr=online_ocr,
            ocr_backend=ocr_backend,
            ocr_max_pages=max_pages,
            scan_detect_sample_pages=8,
            ocr_timeout_seconds=int(os.getenv("PDF_OCR_TIMEOUT_SECONDS", "60")),
            ocr_retries=int(os.getenv("PDF_OCR_RETRIES", "1")),
        )
        parsed = parse_document_result(file_path, cfg)
        return self._analyze_pages(parsed.pages, parsed.metadata)

    def _validate_parse_result(self, result: Dict[str, Any], ext: str) -> Dict[str, Any]:
        """避免空解析被标记为成功，给前端清晰的修复方向。"""
        if result.get("error"):
            return result
        total_chars = int(result.get("total_chars") or 0)
        page_count = int(result.get("pages") or 0)
        parsed_pages = int(result.get("parsed_pages") or 0)
        if total_chars >= 20:
            return result

        if ext == ".pdf" and page_count > 0 and parsed_pages == 0:
            return self._empty_result(
                error=(
                    "PDF 没有提取到可用文字。它很可能是扫描版教材，需要可用 OCR。"
                    "请配置 PDF_OCR_BACKEND=qwen_vl 并设置 MULTIMODAL_API_KEY，"
                    "或安装 PaddleOCR/Tesseract 后重新解析。"
                ),
                pages=page_count,
                failed_pages=result.get("failed_pages", 0),
            )
        if ext in {".ppt", ".doc"}:
            return self._empty_result(
                error=f"{ext} 是旧版二进制格式，当前环境未能成功转换解析。请另存为 {'PPTX' if ext == '.ppt' else 'DOCX'} 后上传，或安装 LibreOffice 后重试。"
            )
        return self._empty_result(error="文件内容为空或未提取到有效文本，请检查文件是否可复制文字，或换成新版 Office/PDF 文件后重试。")

    def _parse_builtin(self, file_path: str, ext: str) -> Dict[str, Any]:
        """内置简化解析（兜底）"""
        text = ""
        from_ocr = False

        if ext == ".pdf":
            text = self._extract_pdf(file_path)
        elif ext == ".docx":
            text = self._extract_docx(file_path)
        elif ext == ".doc":
            return self._empty_result(error=".doc 旧版 Word 需要 RAG/LibreOffice 转换支持。请另存为 .docx 后上传。")
        elif ext == ".pptx":
            text = self._extract_pptx(file_path)
        elif ext == ".ppt":
            return self._empty_result(error=".ppt 旧版 PowerPoint 需要 RAG/LibreOffice 转换支持。请另存为 .pptx 后上传。")
        elif ext == ".txt":
            with open(file_path, encoding="utf-8", errors="ignore") as f:
                text = f.read()

        return self._analyze_text(text, from_ocr=from_ocr)

    def _extract_pdf(self, path: str) -> str:
        try:
            import PyMuPDF
            doc = PyMuPDF.open(path)
            pages = [page.get_text() for page in doc]
            doc.close()
            return "\n".join(pages)
        except ImportError:
            pass
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                return "\n".join(p.extract_text() or "" for p in pdf.pages)
        except Exception as e:
            logger.warning(f"PDF 提取失败: {e}")
            return ""

    def _extract_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.warning(f"DOCX 提取失败: {e}")
            return ""

    def _extract_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"PPTX 提取失败: {e}")
            return ""

    def _analyze_text(self, text: str, from_ocr: bool = False) -> Dict[str, Any]:
        """从文本中分析章节/知识点等结构"""
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        total_chars = len(text)

        # 章节识别（简化规则）
        chapter_patterns = [
            r'^#{1,6}\s+',
            r'^第[一二三四五六七八九十百\d]+章',
            r'^Chapter\s+\d+',
            r'^\d+\.\s+[^\d]',
        ]
        chapters = []
        for line in lines:
            for pat in chapter_patterns:
                if re.match(pat, line):
                    chapters.append(self._clean_heading(line)[:80])
                    break
        chapters = self._dedupe_keep_order(chapters)

        # 公式识别（含数学符号）
        formula_keywords = ['=', '∫', '∑', '∏', 'lim', 'sin', 'cos', 'tan', 'log', '→', '∞', 'Ω', 'V', 'A']
        formulas = [l for l in lines if any(k in l for k in formula_keywords) and len(l) < 200]

        # 例题识别
        examples = [l for l in lines if re.match(r'^(例|例题)\d*[：:]|^Example\s+\d+', l)]

        # 习题识别
        exercises = [
            l for l in lines
            if re.match(r'^(习题|练习|作业|题目|求|Exercise)\s*\d*[：:]?', l) or re.match(r'^\d+[\.、]\s*', l)
        ]

        # 知识点（定义/定理/性质）
        kp_lines = [l for l in lines if re.match(r'^(定义|定理|性质|引理|推论|命题|公式|方法|结论)\s*\d*[：:]', l)]
        kp_lines.extend(self._infer_knowledge_points(lines, formulas))
        kp_lines = self._dedupe_keep_order([self._clean_kp_name(kp) for kp in kp_lines if kp.strip()])

        if not chapters:
            chapters = self._fallback_chapters(lines, kp_lines, formulas, exercises)

        return {
            "text": text,
            "total_chars": total_chars,
            "pages": max(1, total_chars // 2000),   # 估算页数
            "chapters": chapters,
            "knowledge_points": kp_lines[:50],
            "formulas": formulas[:100],
            "examples": examples[:50],
            "exercises": exercises[:50],
            "from_ocr": from_ocr,
            "error": "",
        }

    def _analyze_pages(self, pages: List[Dict[str, Any]], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """面向教材 PDF 的页级结构化分析。"""
        page_items = []
        for page in pages:
            text = page.get("text", "") or ""
            page_no = int(page.get("page_number") or page.get("page") or 0)
            if text.strip():
                page_items.append(
                    {
                        "page": page_no,
                        "text": text,
                        "is_ocr": bool(page.get("is_ocr") or page.get("is_ocr_text")),
                        "method": page.get("extraction_method") or page.get("method", ""),
                    }
                )

        text = "\n\n".join(item["text"] for item in page_items)
        if not page_items:
            result = self._analyze_text(text, from_ocr=False)
            result.update(
                {
                    "pages": metadata.get("page_count", 0),
                    "parsed_pages": 0,
                    "ocr_pages": 0,
                    "failed_pages": metadata.get("failed_page_count", 0),
                    "page_chunks": [],
                    "chapter_sections": [],
                    "knowledge_points_structured": [],
                }
            )
            return result

        line_items: List[Dict[str, Any]] = []
        for item in page_items:
            for raw in item["text"].splitlines():
                line = self._normalize_line(raw)
                if line:
                    line_items.append({"text": line, "page": item["page"], "is_ocr": item["is_ocr"]})

        chapters = self._extract_chapter_sections(line_items, text)
        self._annotate_line_chapters(line_items, chapters)
        structured_kps = self._extract_structured_kps(line_items, chapters, text)
        formulas = self._extract_formulas(line_items)
        examples = self._extract_examples(line_items)
        exercises = self._extract_exercises(line_items)
        chunks = self._build_page_chunks(page_items, chapters)

        chapter_names = [c["name"] for c in chapters]
        if not chapter_names:
            fallback = self._analyze_text(text, from_ocr=any(p["is_ocr"] for p in page_items))
            chapter_names = fallback["chapters"]
            chapters = [
                {"name": name, "number": str(i + 1), "page_start": 0, "intro": f"由教材内容自动归纳：{name}"}
                for i, name in enumerate(chapter_names)
            ]
        kp_names = [kp["name"] for kp in structured_kps]
        if not kp_names:
            kp_names = self._infer_knowledge_points([x["text"] for x in line_items], formulas)
            structured_kps = [
                {
                    "name": name,
                    "definition": name,
                    "kp_type": self._guess_kp_type_rich(name),
                    "difficulty": "medium",
                    "page": "",
                    "chapter_index": i % max(len(chapters), 1),
                }
                for i, name in enumerate(kp_names)
            ]

        return {
            "text": text,
            "total_chars": len(text),
            "pages": int(metadata.get("page_count") or len(pages)),
            "parsed_pages": len(page_items),
            "ocr_pages": sum(1 for p in page_items if p["is_ocr"]),
            "failed_pages": int(metadata.get("failed_page_count") or 0),
            "chapters": chapter_names[:80],
            "chapter_sections": chapters[:80],
            "knowledge_points": kp_names[:300],
            "knowledge_points_structured": structured_kps[:300],
            "formulas": formulas[:300],
            "examples": examples[:200],
            "exercises": exercises[:200],
            "page_chunks": chunks,
            "from_ocr": any(p["is_ocr"] for p in page_items),
            "error": "",
        }

    def _extract_chapter_sections(self, line_items: List[Dict[str, Any]], full_text: str) -> List[Dict[str, Any]]:
        chapters: List[Dict[str, Any]] = []
        seen = set()
        for item in line_items:
            line = self._clean_heading(item["text"])
            title = self._clean_toc_title(line)
            if not self._looks_like_heading(title):
                continue
            key = re.sub(r"\s+", "", title)
            if key in seen:
                continue
            seen.add(key)
            number = self._extract_heading_number(title) or str(len(chapters) + 1)
            chapters.append(
                {
                    "name": title[:120],
                    "number": number[:32],
                    "page_start": item["page"],
                    "intro": f"教材第 {item['page']} 页附近识别到：{title[:90]}",
                }
            )
            if len(chapters) >= 80:
                break

        if chapters:
            for idx, ch in enumerate(chapters):
                if idx + 1 < len(chapters):
                    ch["page_end"] = max(ch["page_start"], chapters[idx + 1]["page_start"] - 1)
                else:
                    ch["page_end"] = 0
            return chapters

        toc_chapters = self._extract_toc_chapters(full_text)
        return toc_chapters[:80]

    def _extract_toc_chapters(self, text: str) -> List[Dict[str, Any]]:
        chapters = []
        seen = set()
        for raw in text.splitlines():
            line = self._clean_toc_title(self._normalize_line(raw))
            if not line or not self._looks_like_heading(line):
                continue
            key = re.sub(r"\s+", "", line)
            if key in seen:
                continue
            seen.add(key)
            chapters.append(
                {
                    "name": line[:120],
                    "number": self._extract_heading_number(line) or str(len(chapters) + 1),
                    "page_start": 0,
                    "page_end": 0,
                    "intro": f"由目录或正文标题识别：{line[:90]}",
                }
            )
        return chapters

    def _extract_structured_kps(
        self,
        line_items: List[Dict[str, Any]],
        chapters: List[Dict[str, Any]],
        full_text: str,
    ) -> List[Dict[str, Any]]:
        candidates: List[Dict[str, Any]] = []
        explicit = re.compile(r"^(定义|定理|性质|引理|推论|命题|公式|法则|例题|例)\s*([0-9一二三四五六七八九十.．-]*)\s*[：:]?\s*(.+)?$")
        for idx, item in enumerate(line_items):
            line = item["text"]
            match = explicit.match(line)
            if match:
                prefix = match.group(1)
                tail = (match.group(3) or "").strip()
                context = self._join_nearby_lines(line_items, idx, max_lines=3)
                name = self._clean_kp_name(f"{prefix}{match.group(2)} {tail}".strip()) or prefix
                candidates.append(
                    self._kp_record(
                        name=name,
                        definition=context,
                        kp_type=self._guess_kp_type_rich(prefix),
                        page=item["page"],
                        chapter_index=int(item.get("chapter_index", self._chapter_index_for_page(chapters, item["page"]))),
                    )
                )
                continue
            if self._looks_like_section_title(line):
                candidates.append(
                    self._kp_record(
                        name=self._clean_toc_title(line),
                        definition=self._join_nearby_lines(line_items, idx, max_lines=2),
                        kp_type="concept",
                        page=item["page"],
                        chapter_index=int(item.get("chapter_index", self._chapter_index_for_page(chapters, item["page"]))),
                    )
                )

        domain_terms = [
            "随机试验", "样本空间", "随机事件", "事件关系", "事件运算", "频率", "概率", "古典概型", "几何概型",
            "条件概率", "乘法公式", "全概率公式", "贝叶斯公式", "事件独立性", "伯努利试验",
            "随机变量", "分布函数", "离散型随机变量", "连续型随机变量", "概率密度", "分布律",
            "数学期望", "方差", "协方差", "相关系数", "矩母函数", "大数定律", "中心极限定理",
            "正态分布", "二项分布", "泊松分布", "指数分布", "均匀分布", "卡方分布", "t分布", "F分布",
            "参数估计", "最大似然估计", "矩估计", "置信区间", "假设检验", "显著性水平", "回归分析",
        ]
        seen_text = full_text
        for term in domain_terms:
            if term in seen_text:
                page = self._first_page_containing(line_items, term)
                chapter_index = self._first_chapter_index_containing(line_items, term, chapters)
                candidates.append(
                    self._kp_record(
                        name=term,
                        definition=self._context_for_term(line_items, term),
                        kp_type="concept",
                        page=page,
                        chapter_index=chapter_index,
                    )
                )

        formulas = self._extract_formulas(line_items)
        for formula in formulas[:80]:
            text = str(formula)
            source_line = self._first_line_equal(line_items, text)
            page = int(source_line.get("page", 0)) if source_line else 0
            candidates.append(
                self._kp_record(
                    name=f"公式：{text[:80]}",
                    definition=text,
                    kp_type="formula",
                    page=page,
                    chapter_index=int(source_line.get("chapter_index", self._chapter_index_for_page(chapters, page))) if source_line else self._chapter_index_for_page(chapters, page),
                )
            )

        return self._dedupe_kp_records(candidates)

    def _extract_formulas(self, line_items: List[Dict[str, Any]]) -> List[str]:
        keywords = ["=", "∫", "∑", "∏", "lim", "sin", "cos", "tan", "log", "ln", "→", "∞", "P{", "P(", "E(", "D(", "Var", "Cov"]
        formulas = []
        for item in line_items:
            line = item["text"]
            if 4 <= len(line) <= 220 and any(k in line for k in keywords):
                formulas.append(line)
        return self._dedupe_keep_order(formulas)

    def _extract_examples(self, line_items: List[Dict[str, Any]]) -> List[str]:
        return self._dedupe_keep_order(
            [item["text"] for item in line_items if re.match(r"^(例|例题)\s*[0-9一二三四五六七八九十.．-]*", item["text"])]
        )

    def _extract_exercises(self, line_items: List[Dict[str, Any]]) -> List[str]:
        result = []
        for item in line_items:
            line = item["text"]
            if re.match(r"^(习题|练习|总习题|复习题|思考题)\s*[0-9一二三四五六七八九十.．-]*", line):
                result.append(line)
            elif re.match(r"^[0-9]+[\.．、]\s*", line) and any(word in line for word in ["证明", "求", "计算", "设", "试"]):
                result.append(line)
        return self._dedupe_keep_order(result)

    def _build_page_chunks(self, page_items: List[Dict[str, Any]], chapters: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        chunks: List[Dict[str, Any]] = []
        chunk_size = 900
        overlap = 120
        for page in page_items:
            text = page["text"].strip()
            if not text:
                continue
            section_spans = self._section_spans_for_page(text, page["page"], chapters)
            for section_text, chapter_index in section_spans:
                start = 0
                while start < len(section_text):
                    chunk = self._clean_chunk_text(section_text[start:start + chunk_size])
                    if len(chunk) >= 10:
                        chunks.append(
                            {
                                "text": chunk,
                                "page_start": page["page"],
                                "page_end": page["page"],
                                "chapter_index": chapter_index,
                                "from_ocr": page["is_ocr"],
                            }
                        )
                    if start + chunk_size >= len(section_text):
                        break
                    start += chunk_size - overlap
        return chunks

    def _section_spans_for_page(self, text: str, page: int, chapters: List[Dict[str, Any]]) -> List[Tuple[str, int]]:
        hits = []
        for idx, chapter in enumerate(chapters):
            if int(chapter.get("page_start") or 0) != page:
                continue
            name = str(chapter.get("name") or "").strip()
            if not name:
                continue
            pos = text.find(name)
            if pos >= 0:
                hits.append((pos, idx))
        hits = sorted(set(hits), key=lambda item: item[0])
        if not hits:
            return [(text, self._chapter_index_for_page(chapters, page))]
        spans: List[Tuple[str, int]] = []
        for i, (pos, chapter_index) in enumerate(hits):
            end = hits[i + 1][0] if i + 1 < len(hits) else len(text)
            part = text[pos:end].strip()
            if part:
                spans.append((part, chapter_index))
        prefix = text[:hits[0][0]].strip()
        if prefix and len(prefix) >= 40:
            spans.insert(0, (prefix, max(0, hits[0][1] - 1)))
        return spans or [(text, self._chapter_index_for_page(chapters, page))]

    def _clean_chunk_text(self, value: str) -> str:
        text = value.strip()
        text = re.sub(r"\n?\s*#{1,6}\s*$", "", text).strip()
        text = re.sub(r"^\s*#{1,6}\s*", "", text).strip()
        return text

    def _annotate_line_chapters(self, line_items: List[Dict[str, Any]], chapters: List[Dict[str, Any]]) -> None:
        chapter_by_name = {str(ch.get("name") or ""): idx for idx, ch in enumerate(chapters)}
        current = 0
        for item in line_items:
            line = self._clean_toc_title(item["text"])
            if line in chapter_by_name:
                current = chapter_by_name[line]
            else:
                current = self._chapter_index_for_page(chapters, item["page"]) if not chapter_by_name else current
            item["chapter_index"] = current

    def _first_chapter_index_containing(self, line_items: List[Dict[str, Any]], term: str, chapters: List[Dict[str, Any]]) -> int:
        for item in line_items:
            if term in item["text"]:
                return int(item.get("chapter_index", self._chapter_index_for_page(chapters, item["page"])))
        return 0

    def _normalize_line(self, value: str) -> str:
        text = value.strip()
        text = re.sub(r"\s+", " ", text)
        text = text.replace("．", ".").replace("：", ":")
        return text.strip()

    def _clean_toc_title(self, value: str) -> str:
        text = self._clean_heading(value)
        text = re.sub(r"[.·•…．]{2,}\s*\d+\s*$", "", text)
        text = re.sub(r"\s+\d{1,4}\s*$", "", text)
        text = re.sub(r"^\d+\s+", "", text) if text.startswith("目录 ") else text
        return text.strip(" -—\t")

    def _looks_like_heading(self, line: str) -> bool:
        if not line or len(line) > 90:
            return False
        patterns = [
            r"^第[一二三四五六七八九十百零〇0-9]+[章节篇]\s*.+",
            r"^[0-9]+(?:\.[0-9]+){1,3}\s+[\u4e00-\u9fa5A-Za-z].+",
            r"^§\s*[0-9]+(?:\.[0-9]+)*\s*.+",
            r"^Chapter\s+[0-9IVXLC]+.+",
        ]
        if not any(re.match(pattern, line, re.IGNORECASE) for pattern in patterns):
            return False
        noise = ["如下", "因为", "所以", "于是", "可得", "证明", "解:", "例", "习题"]
        return not any(line.startswith(x) for x in noise)

    def _looks_like_section_title(self, line: str) -> bool:
        if not self._looks_like_heading(line):
            return False
        return bool(re.match(r"^[0-9]+(?:\.[0-9]+){1,3}\s+", line) or re.match(r"^第[一二三四五六七八九十百零〇0-9]+节", line))

    def _extract_heading_number(self, line: str) -> str:
        match = re.match(r"^(第[一二三四五六七八九十百零〇0-9]+[章节篇]|[0-9]+(?:\.[0-9]+){0,3}|§\s*[0-9]+(?:\.[0-9]+)*)", line)
        return match.group(1).strip() if match else ""

    def _join_nearby_lines(self, line_items: List[Dict[str, Any]], idx: int, max_lines: int = 3) -> str:
        page = line_items[idx]["page"]
        lines = []
        for item in line_items[idx:idx + max_lines]:
            if item["page"] != page:
                break
            lines.append(item["text"])
        return " ".join(lines)[:500]

    def _first_page_containing(self, line_items: List[Dict[str, Any]], term: str) -> int:
        for item in line_items:
            if term in item["text"]:
                return int(item["page"])
        return 0

    def _first_line_equal(self, line_items: List[Dict[str, Any]], text: str) -> Optional[Dict[str, Any]]:
        for item in line_items:
            if item["text"] == text:
                return item
        return None

    def _context_for_term(self, line_items: List[Dict[str, Any]], term: str) -> str:
        for idx, item in enumerate(line_items):
            if term in item["text"]:
                start = max(0, idx - 1)
                end = min(len(line_items), idx + 3)
                same_page = [x["text"] for x in line_items[start:end] if x["page"] == item["page"]]
                return " ".join(same_page)[:500]
        return term

    def _chapter_index_for_page(self, chapters: List[Dict[str, Any]], page: int) -> int:
        if not chapters or not page:
            return 0
        chosen = 0
        for idx, ch in enumerate(chapters):
            start = int(ch.get("page_start") or 0)
            if start and start <= page:
                chosen = idx
            elif start and start > page:
                break
        return chosen

    def _kp_record(
        self,
        name: str,
        definition: str,
        kp_type: str,
        page: int,
        chapter_index: int,
    ) -> Dict[str, Any]:
        clean_name = self._clean_kp_name(name)
        return {
            "name": clean_name[:120],
            "definition": definition[:800] if definition else clean_name,
            "kp_type": kp_type,
            "difficulty": self._guess_difficulty(clean_name, definition),
            "page": str(page) if page else "",
            "chapter_index": max(0, chapter_index),
        }

    def _dedupe_kp_records(self, items: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        seen = set()
        result = []
        for item in items:
            name = item.get("name", "").strip()
            if not name or len(name) < 2:
                continue
            key = re.sub(r"\s+", "", name)
            if key in seen:
                continue
            seen.add(key)
            result.append(item)
        return result

    def _guess_kp_type_rich(self, text: str) -> str:
        if any(word in text for word in ["定理", "引理", "推论", "法则", "律"]):
            return "theorem"
        if any(word in text for word in ["公式", "=", "∫", "∑", "P(", "P{", "E(", "D("]):
            return "formula"
        if any(word in text for word in ["例", "题", "解法", "方法", "估计", "检验"]):
            return "method"
        return "concept"

    def _guess_difficulty(self, name: str, definition: str) -> str:
        hard_words = ["中心极限定理", "最大似然", "假设检验", "协方差", "条件概率", "贝叶斯", "矩母函数"]
        if any(word in name or word in definition for word in hard_words):
            return "hard"
        easy_words = ["样本空间", "随机事件", "概率", "频率", "分布函数"]
        if any(word == name for word in easy_words):
            return "easy"
        return "medium"

    def _fallback_chapters(self, lines: List[str], kps: List[str], formulas: List[str], exercises: List[str]) -> List[str]:
        """当原文没有显式章节标题时，生成可用于学习路径的章节。"""
        if not lines:
            return ["文档概览"]
        chapters = []
        title = self._clean_heading(lines[0])[:40]
        if 2 <= len(title) <= 40 and not any(mark in title for mark in ["=", "：", ":"]):
            chapters.append(title)
        chapters.append("核心内容")
        if formulas:
            chapters.append("公式与计算方法")
        if exercises:
            chapters.append("例题与练习")
        if len(kps) >= 6:
            chapters.append("重点知识点梳理")
        return self._dedupe_keep_order(chapters)[:6]

    def _infer_knowledge_points(self, lines: List[str], formulas: List[str]) -> List[str]:
        """从普通文本中兜底提取知识点，避免 TXT/PPT 无结构时为 0。"""
        import re
        joined = "\n".join(lines)
        candidates: List[str] = []
        domain_terms = [
            "电动势", "内阻", "闭合电路欧姆定律", "串联电路", "并联电路", "外电路总电阻", "路端电压", "电功率",
            "极限", "导数", "积分", "函数", "等价无穷小", "链式法则",
            "受力分析", "牛顿第二定律", "摩擦力", "支持力",
            "化学键", "离子键", "共价键", "金属键", "杂化轨道", "分子间作用力", "范德华力", "氢键",
            "数组", "链表", "栈", "队列", "树", "图", "排序", "查找", "递归",
        ]
        for term in domain_terms:
            if term in joined:
                candidates.append(term)
        for line in formulas[:20]:
            compact = re.sub(r"\s+", " ", line).strip()
            if 4 <= len(compact) <= 120:
                candidates.append(f"公式：{compact}")
        for line in lines[:80]:
            if line.startswith("#"):
                continue
            if len(line) > 120:
                continue
            if any(word in line for word in ["已知", "求", "计算", "公式", "定律", "电阻", "电流", "电压", "功率"]):
                candidates.append(line)
        if not candidates and lines:
            for line in lines[:8]:
                if line.startswith("#"):
                    continue
                if 6 <= len(line) <= 80:
                    candidates.append(self._clean_heading(line))
        return candidates[:20]

    @staticmethod
    def _clean_heading(value: str) -> str:
        import re
        return re.sub(r"^#{1,6}\s*", "", value.strip()).strip()

    @staticmethod
    def _clean_kp_name(value: str) -> str:
        import re
        text = re.sub(r"^(定义|定理|性质|引理|推论|命题|公式|方法|结论)\s*\d*[：:]\s*", "", value.strip())
        return text[:120]

    @staticmethod
    def _dedupe_keep_order(items: List[str]) -> List[str]:
        seen = set()
        result = []
        for item in items:
            key = item.strip()
            if not key or key in seen:
                continue
            seen.add(key)
            result.append(key)
        return result

    @staticmethod
    def _empty_result(error: str = "", pages: int = 0, failed_pages: int = 0) -> Dict[str, Any]:
        return {
            "text": "", "total_chars": 0, "pages": pages,
            "parsed_pages": 0, "ocr_pages": 0, "failed_pages": failed_pages,
            "chapters": [], "knowledge_points": [],
            "formulas": [], "examples": [], "exercises": [],
            "chapter_sections": [], "knowledge_points_structured": [], "page_chunks": [],
            "from_ocr": False, "error": error,
        }


# ════════════════════════════════════════════════════════════
#  向量检索服务
# ════════════════════════════════════════════════════════════

class RAGRetrieveService:
    """
    RAG 检索服务
    对 1号 IndexManager / VectorStore 的封装
    支持按 course_id 检索相关文本片段
    """

    def __init__(self):
        self._index_managers: Dict[str, Any] = {}

    def _get_manager(self, course_id: str):
        """懒加载每门课程的索引管理器"""
        if course_id in self._index_managers:
            return self._index_managers[course_id]

        if not _rag_available:
            return None

        try:
            from rag.retrieval.embeddings import EmbeddingProvider  # type: ignore
            from rag.storage.vector_store import VectorStore        # type: ignore
            from rag.pipeline.index_manager import IndexManager     # type: ignore

            emb = EmbeddingProvider(
                provider="local",
                model_name="paraphrase-multilingual-MiniLM-L12-v2",
                cache_dir=str(settings.EMBED_CACHE_DIR),
            )

            from rag.config import VectorStoreConfig  # type: ignore  (config 仍在根目录)
            store = VectorStore(config=VectorStoreConfig(
                collection_name=f"course_{course_id}",
                persist_dir=settings.CHROMA_PERSIST_DIR,
            ))
            manager = IndexManager(store, emb, course_id=course_id)
            self._index_managers[course_id] = manager
            return manager
        except Exception as e:
            logger.warning(f"IndexManager 初始化失败: {e}")
            return None

    def index_document(self, course_id: str, file_path: str, force: bool = False) -> int:
        """将文档索引到向量库，返回 chunk 数量"""
        manager = self._get_manager(course_id)
        if manager:
            try:
                return manager.index_document(file_path, force=force)
            except Exception as e:
                logger.error(f"向量索引失败: {e}")
        return 0

    def retrieve(
        self,
        course_id: str,
        query: str,
        top_k: int = 5,
        chapter: Optional[str] = None,
    ) -> List[str]:
        """
        检索相关片段，返回文本列表

        优先用 1号模块；降级返回空列表（由调用方从 DB 里取 fragments）
        """
        manager = self._get_manager(course_id)
        if manager:
            try:
                # 1号的 VectorStore 检索接口
                store = manager.store
                results = store.query(query_text=query, top_k=top_k)
                return [doc for doc in results.get("documents", [[]])[0]]
            except Exception as e:
                logger.warning(f"向量检索失败: {e}")

        # 降级：从数据库 fragments 表做关键词匹配
        return self._keyword_fallback(course_id, query, top_k)

    def _keyword_fallback(self, course_id: str, query: str, top_k: int) -> List[str]:
        """关键词降级检索（从 SQLite）"""
        try:
            from ..core.database import SessionLocal, Fragment
            db = SessionLocal()
            frags = (
                db.query(Fragment)
                .filter(Fragment.course_id == course_id)
                .limit(top_k * 5)
                .all()
            )
            db.close()

            keywords = query.lower().split()
            scored = []
            for f in frags:
                score = sum(1 for kw in keywords if kw in f.text.lower())
                if score > 0:
                    scored.append((score, f.text))
            scored.sort(reverse=True)
            return [t for _, t in scored[:top_k]]
        except Exception as e:
            logger.warning(f"关键词降级检索失败: {e}")
            return []


# ─── 单例 ────────────────────────────────────────────────────
parse_service = DocumentParseService()
rag_service = RAGRetrieveService()
